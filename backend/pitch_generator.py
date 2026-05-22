import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── Color palette ─────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x0F, 0x1A)
CARD     = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT   = RGBColor(0x63, 0x66, 0xF1)
ACCENT2  = RGBColor(0x8B, 0x5C, 0xF6)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
SUCCESS  = RGBColor(0x10, 0xB9, 0x81)
WARNING  = RGBColor(0xF5, 0x9E, 0x0B)
DANGER   = RGBColor(0xEF, 0x44, 0x44)


def _set_bg(slide, prs):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _add_textbox(slide, text, left, top, width, height,
                  font_size=18, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_bullet_box(slide, items, left, top, width, height,
                     title=None, title_color=ACCENT, item_color=WHITE):
    if title:
        _add_textbox(slide, title, left, top, width, 0.4,
                     font_size=13, bold=True, color=title_color)
        top += 0.45

    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(10)
        run.font.color.rgb = item_color


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_title(prs, brand_name, startup_idea):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, prs)
    _add_rect(slide, 0, 2.5, 10, 0.05, ACCENT)
    _add_textbox(slide, brand_name.upper(), 0.5, 0.8, 9, 1,
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, startup_idea, 0.5, 2.0, 9, 0.8,
                 font_size=18, color=MUTED, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "Startup Analysis & Pitch Deck", 0.5, 2.8, 9, 0.6,
                 font_size=14, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "Confidential — For Discussion Purposes Only",
                 0.5, 6.8, 9, 0.4, font_size=10, color=MUTED, align=PP_ALIGN.CENTER)


def _slide_section_header(prs, number, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, prs)
    _add_rect(slide, 0, 0, 10, 7.5, CARD)
    _add_rect(slide, 0, 0, 0.15, 7.5, ACCENT)
    _add_textbox(slide, f"{number:02d}", 0.5, 1.5, 2, 1.5,
                 font_size=72, bold=True, color=ACCENT)
    _add_textbox(slide, title.upper(), 0.5, 3.2, 9, 0.8,
                 font_size=32, bold=True, color=WHITE)
    _add_textbox(slide, subtitle, 0.5, 4.1, 9, 0.6,
                 font_size=16, color=MUTED)


def _slide_market(prs, market):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, prs)
    _add_textbox(slide, "MARKET ANALYSIS", 0.3, 0.2, 9, 0.5,
                 font_size=22, bold=True, color=WHITE)
    _add_rect(slide, 0.3, 0.75, 9.4, 0.04, ACCENT)

    # TAM SAM SOM boxes
    tam = market.get("tam", {}).get("value", 0)
    sam = market.get("sam", {}).get("value", 0)
    som = market.get("som", {}).get("value", 0)
    cagr = market.get("cagr", 0)

    for i, (label, val, col) in enumerate([
        ("TAM", f"${tam}B", ACCENT),
        ("SAM", f"${sam}B", ACCENT2),
        ("SOM", f"${som}B", SUCCESS),
    ]):
        x = 0.3 + i * 3.2
        _add_rect(slide, x, 1.0, 2.9, 1.6, CARD)
        _add_textbox(slide, label, x + 0.15, 1.1, 2.6, 0.4,
                     font_size=12, bold=True, color=MUTED)
        _add_textbox(slide, val, x + 0.15, 1.45, 2.6, 0.7,
                     font_size=28, bold=True, color=col)

    _add_rect(slide, 0.3, 2.8, 9.4, 0.04, CARD)

    # CAGR + Overview
    _add_textbox(slide, f"CAGR: {cagr}%", 0.3, 2.95, 3, 0.5,
                 font_size=18, bold=True, color=WARNING)
    _add_textbox(slide, market.get("overview", ""), 0.3, 3.55, 9.4, 1.0,
                 font_size=11, color=MUTED)

    # Segments
    segs = market.get("segments", [])
    _add_textbox(slide, "KEY SEGMENTS", 0.3, 4.7, 9, 0.35,
                 font_size=13, bold=True, color=ACCENT)
    for i, seg in enumerate(segs[:3]):
        x = 0.3 + i * 3.2
        _add_rect(slide, x, 5.1, 2.9, 1.8, CARD)
        _add_textbox(slide, seg.get("name", ""), x + 0.15, 5.2, 2.6, 0.35,
                     font_size=11, bold=True, color=WHITE)
        _add_textbox(slide, seg.get("size", ""), x + 0.15, 5.55, 2.6, 0.3,
                     font_size=10, color=SUCCESS)
        pain = " • ".join(seg.get("pain_points", [])[:2])
        _add_textbox(slide, pain, x + 0.15, 5.85, 2.6, 0.9,
                     font_size=9, color=MUTED)


def _slide_competitors(prs, comp):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, prs)
    _add_textbox(slide, "COMPETITIVE LANDSCAPE", 0.3, 0.2, 9, 0.5,
                 font_size=22, bold=True, color=WHITE)
    _add_rect(slide, 0.3, 0.75, 9.4, 0.04, ACCENT)

    competitors = comp.get("competitors", [])[:5]
    for i, c in enumerate(competitors):
        y = 1.0 + i * 1.2
        _add_rect(slide, 0.3, y, 9.4, 1.1, CARD)
        _add_textbox(slide, c.get("name", ""), 0.5, y + 0.05, 2, 0.4,
                     font_size=13, bold=True, color=WHITE)
        _add_textbox(slide, c.get("funding", ""), 0.5, y + 0.5, 1.5, 0.35,
                     font_size=10, color=SUCCESS)
        _add_textbox(slide, c.get("pricing", ""), 2.2, y + 0.05, 2, 0.35,
                     font_size=10, color=WARNING)
        usps = " | ".join(c.get("usps", [])[:2])
        _add_textbox(slide, f"✓ {usps}", 2.2, y + 0.45, 4, 0.35,
                     font_size=9, color=MUTED)
        weaknesses = " | ".join(c.get("weaknesses", [])[:2])
        _add_textbox(slide, f"✗ {weaknesses}", 6.3, y + 0.05, 3.2, 0.8,
                     font_size=9, color=DANGER)

    gaps = comp.get("gaps", [])
    if gaps:
        _add_textbox(slide, "MARKET GAPS WE ADDRESS", 0.3, 7.05, 9, 0.3,
                     font_size=11, bold=True, color=ACCENT)


def _slide_funding(prs, funding):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, prs)
    _add_textbox(slide, "FUNDING LANDSCAPE", 0.3, 0.2, 9, 0.5,
                 font_size=22, bold=True, color=WHITE)
    _add_rect(slide, 0.3, 0.75, 9.4, 0.04, ACCENT)

    sentiment = funding.get("sentiment", "neutral").upper()
    total = funding.get("total_investment", "N/A")
    sent_color = SUCCESS if sentiment == "BULLISH" else (DANGER if sentiment == "BEARISH" else WARNING)

    _add_textbox(slide, f"Investor Sentiment: {sentiment}", 0.3, 0.9, 5, 0.5,
                 font_size=16, bold=True, color=sent_color)
    _add_textbox(slide, f"Total Investment: {total}", 5.5, 0.9, 4, 0.5,
                 font_size=16, bold=True, color=ACCENT)

    # Recent rounds table header
    _add_textbox(slide, "RECENT FUNDING ROUNDS", 0.3, 1.55, 9, 0.35,
                 font_size=12, bold=True, color=ACCENT)
    _add_rect(slide, 0.3, 1.95, 9.4, 0.04, ACCENT)

    rounds = funding.get("rounds", [])[:5]
    for i, r in enumerate(rounds):
        y = 2.05 + i * 0.62
        bg = CARD if i % 2 == 0 else BG
        _add_rect(slide, 0.3, y, 9.4, 0.58, bg)
        _add_textbox(slide, r.get("company", ""), 0.45, y + 0.1, 2.2, 0.38,
                     font_size=11, bold=True, color=WHITE)
        _add_textbox(slide, r.get("stage", ""), 2.7, y + 0.1, 1.5, 0.38,
                     font_size=10, color=ACCENT2)
        _add_textbox(slide, r.get("amount", ""), 4.3, y + 0.1, 1.5, 0.38,
                     font_size=11, bold=True, color=SUCCESS)
        _add_textbox(slide, r.get("investor", ""), 5.9, y + 0.1, 2.5, 0.38,
                     font_size=10, color=MUTED)
        _add_textbox(slide, str(r.get("year", "")), 8.5, y + 0.1, 1, 0.38,
                     font_size=10, color=MUTED)

    rec = funding.get("recommendations", {})
    _add_textbox(slide,
                 f"Recommendation: {rec.get('ideal_stage','')} — {rec.get('round_size','')}",
                 0.3, 5.25, 9.4, 0.45, font_size=13, bold=True, color=WARNING)


def _slide_swot(prs, swot):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, prs)
    _add_textbox(slide, "SWOT ANALYSIS", 0.3, 0.2, 9, 0.5,
                 font_size=22, bold=True, color=WHITE)
    _add_rect(slide, 0.3, 0.75, 9.4, 0.04, ACCENT)

    quadrants = [
        ("STRENGTHS",     swot.get("strengths", []),     SUCCESS, 0.3,  1.0),
        ("WEAKNESSES",    swot.get("weaknesses", []),    DANGER,  5.1,  1.0),
        ("OPPORTUNITIES", swot.get("opportunities", []), ACCENT,  0.3,  4.3),
        ("THREATS",       swot.get("threats", []),       WARNING, 5.1,  4.3),
    ]

    for title, items, color, x, y in quadrants:
        _add_rect(slide, x, y, 4.6, 3.1, CARD)
        _add_textbox(slide, title, x + 0.15, y + 0.1, 4.2, 0.4,
                     font_size=12, bold=True, color=color)
        _add_rect(slide, x + 0.15, y + 0.55, 4.2, 0.03, color)
        for j, item in enumerate(items[:4]):
            _add_textbox(slide, f"• {item}", x + 0.15, y + 0.65 + j * 0.55, 4.2, 0.5,
                         font_size=9, color=WHITE)


def _slide_gtm(prs, gtm):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, prs)
    _add_textbox(slide, "GO-TO-MARKET STRATEGY", 0.3, 0.2, 9, 0.5,
                 font_size=22, bold=True, color=WHITE)
    _add_rect(slide, 0.3, 0.75, 9.4, 0.04, ACCENT)

    # Value proposition
    _add_textbox(slide, f'"{gtm.get("value_proposition", "")}"',
                 0.3, 0.9, 9.4, 0.7, font_size=13, italic=True, color=ACCENT)

    # Pricing tiers
    _add_textbox(slide, "PRICING", 0.3, 1.7, 9, 0.35,
                 font_size=12, bold=True, color=ACCENT)
    pricing = gtm.get("pricing", [])[:3]
    tier_colors = [SUCCESS, ACCENT, ACCENT2]
    for i, tier in enumerate(pricing):
        x = 0.3 + i * 3.2
        _add_rect(slide, x, 2.1, 2.9, 1.4, CARD)
        _add_textbox(slide, tier.get("tier", ""), x + 0.15, 2.2, 2.6, 0.35,
                     font_size=12, bold=True, color=tier_colors[i % 3])
        _add_textbox(slide, tier.get("price", ""), x + 0.15, 2.55, 2.6, 0.4,
                     font_size=16, bold=True, color=WHITE)
        incs = " • ".join(tier.get("inclusions", [])[:2])
        _add_textbox(slide, incs, x + 0.15, 2.95, 2.6, 0.45,
                     font_size=8, color=MUTED)

    # 5-Phase roadmap
    _add_textbox(slide, "5-PHASE ROADMAP", 0.3, 3.65, 9, 0.35,
                 font_size=12, bold=True, color=ACCENT)
    phases = gtm.get("phases", [])[:5]
    phase_colors = [SUCCESS, ACCENT, ACCENT2, WARNING, DANGER]
    for i, ph in enumerate(phases):
        x = 0.3 + i * 1.9
        _add_rect(slide, x, 4.05, 1.75, 2.0, CARD)
        _add_rect(slide, x, 4.05, 1.75, 0.3, phase_colors[i % 5])
        _add_textbox(slide, f"P{ph.get('phase', i+1)}", x + 0.05, 4.07, 0.5, 0.25,
                     font_size=9, bold=True, color=BG)
        _add_textbox(slide, ph.get("months", ""), x + 0.05, 4.4, 1.65, 0.3,
                     font_size=8, bold=True, color=phase_colors[i % 5])
        _add_textbox(slide, ph.get("title", ""), x + 0.05, 4.7, 1.65, 0.35,
                     font_size=9, color=WHITE)
        goals = " ".join(ph.get("goals", [])[:1])
        _add_textbox(slide, goals, x + 0.05, 5.1, 1.65, 0.85,
                     font_size=8, color=MUTED)

    # KPIs
    kpis = gtm.get("kpis", {})
    _add_textbox(slide,
                 f"12M Revenue Target: ${kpis.get('revenue_12month', 0):,}  |  "
                 f"MRR (12M): ${kpis.get('mrr_12month', 0):,}  |  "
                 f"CAC: ${kpis.get('cac', 0):,}  |  "
                 f"LTV: ${kpis.get('ltv', 0):,}",
                 0.3, 6.2, 9.4, 0.5, font_size=11, bold=True, color=WARNING)


def _slide_closing(prs, brand_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, prs)
    _add_rect(slide, 0, 0, 10, 7.5, CARD)
    _add_rect(slide, 0, 3.45, 10, 0.1, ACCENT)
    _add_textbox(slide, "LET'S BUILD THE FUTURE", 0.5, 1.5, 9, 1,
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, brand_name, 0.5, 2.8, 9, 0.7,
                 font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "[Your Name]", 0.5, 4.2, 9, 0.5,
                 font_size=18, color=MUTED, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "[your@email.com]  |  [LinkedIn]  |  [Website]",
                 0.5, 4.8, 9, 0.5, font_size=14, color=MUTED, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "Replace placeholder text with your actual contact details",
                 0.5, 6.8, 9, 0.4, font_size=9, italic=True,
                 color=RGBColor(0x55, 0x55, 0x77), align=PP_ALIGN.CENTER)


# ── Public entry point ────────────────────────────────────────────────────────
def generate_pitch_deck(analysis: dict, brand_name: str = "Your Brand") -> str:
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    startup_idea = analysis.get("startup_idea", "Startup")

    _slide_title(prs, brand_name, startup_idea)
    _slide_section_header(prs, 1, "Market Analysis",
                          "TAM · SAM · SOM · CAGR · Customer Segments")
    _slide_market(prs, analysis.get("market", {}))
    _slide_section_header(prs, 2, "Competitive Landscape",
                          "Top Competitors · USPs · Weaknesses · Market Gaps")
    _slide_competitors(prs, analysis.get("competitors", {}))
    _slide_section_header(prs, 3, "Funding Landscape",
                          "Recent Rounds · Active VCs · Investor Metrics")
    _slide_funding(prs, analysis.get("funding", {}))
    _slide_section_header(prs, 4, "SWOT Analysis",
                          "Strengths · Weaknesses · Opportunities · Threats")
    _slide_swot(prs, analysis.get("swot", {}))
    _slide_section_header(prs, 5, "Go-To-Market Strategy",
                          "ICP · Pricing · Channels · Roadmap · KPIs")
    _slide_gtm(prs, analysis.get("gtm", {}))
    _slide_closing(prs, brand_name)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name
